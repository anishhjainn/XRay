# utils/text_extract.py
"""
Lightweight, read-only text extraction helpers for supported file types.

Why this module (SOLID):
- SRP: Only extracts and normalizes text; no UI, no checks policy, no logging side effects.
- OCP: Adding file types or swapping strategies doesn't change orchestrator/checks; add a function here.
- LSP: Processors can switch to alternative extractors without affecting callers (they only expect a string).
- ISP: Small, focused functions; processors import only what they use.
- DIP: Checks depend on plain text in metadata (provided by processors), not on heavy parsing libraries.

Public API:
- normalize_text(s: str) -> str
- tokenize_words(s: str) -> list[str]
- extract_docx_text(path, max_chars) -> str
- extract_pptx_text(path, max_chars) -> str
- extract_xlsx_text(path, max_chars, include_hidden=True, skip_formulas=True) -> str
- extract_pdf_text(path, max_chars) -> str

All extractors enforce the max_chars cap by stopping early.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Iterable, List, Optional, Union

PathLike = Union[str, Path]


# ---------------- Normalization and tokenization ----------------


_WHITESPACE_RE = re.compile(r"\s+", re.UNICODE)
_WORD_RE_EN = re.compile(r"[A-Za-z']+")


def normalize_text(s: str) -> str:
    """
    Normalize extracted text for downstream analysis:
    - Collapse any whitespace runs (tabs/newlines) into single spaces.
    - Strip leading/trailing whitespace.
    - Keep punctuation intact (useful for future grammar heuristics).
    """
    if not s:
        return ""
    s = _WHITESPACE_RE.sub(" ", s)
    return s.strip()


def tokenize_words(s: str) -> List[str]:
    """
    Basic word tokenizer for English-like text:
    - Keeps alphabetic sequences and apostrophes (e.g., don't, it's).
    - Lowercases tokens for stable matching in spelling.
    - Skips numbers and punctuation.
    """
    if not s:
        return []
    return [m.group(0).lower() for m in _WORD_RE_EN.finditer(s)]


# ---------------- Internal helpers ----------------


def _append_and_maybe_stop(chunks: List[str], to_add: str, max_chars: int) -> bool:
    """
    Append 'to_add' to chunks if we haven't reached max_chars yet.
    Returns True if we should stop (limit reached), False otherwise.

    We avoid constantly joining the whole text; we only join once at the end.
    """
    if not to_add:
        return False
    current_len = sum(len(c) for c in chunks)
    remaining = max_chars - current_len
    if remaining <= 0:
        return True  # already at or above cap
    if len(to_add) <= remaining:
        chunks.append(to_add)
        return False
    # Need only part of it
    chunks.append(to_add[:remaining])
    return True


def _safe_str(x) -> str:
    return "" if x is None else str(x)


# ---------------- File type extractors ----------------


def extract_docx_text(path: PathLike, max_chars: int) -> str:
    """
    Extract plain text from a .docx file:
    - Gather paragraph text.
    - Gather table cell text.
    - Stop early once 'max_chars' is reached.
    """
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("python-docx is required to extract DOCX text") from exc

    p = Path(path)
    chunks: List[str] = []
    try:
        doc = Document(str(p))

        # Paragraphs
        for para in doc.paragraphs:
            if _append_and_maybe_stop(chunks, _safe_str(para.text) + "\n", max_chars):
                return normalize_text("".join(chunks))

        # Tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if _append_and_maybe_stop(chunks, _safe_str(cell.text) + "\n", max_chars):
                        return normalize_text("".join(chunks))

    except Exception as exc:
        # Propagate: processors catch and set metadata["read_error"]
        raise

    return normalize_text("".join(chunks))


def extract_pptx_text(path: PathLike, max_chars: int) -> str:
    """
    Extract plain text from a .pptx file:
    - For each slide, collect:
      - Text from shapes with text_frame (all paragraphs/runs).
      - Text from table cells.
    - Stop early once 'max_chars' is reached.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to extract PPTX text") from exc

    p = Path(path)
    chunks: List[str] = []
    try:
        prs = Presentation(str(p))
        for slide in prs.slides:
            for shape in slide.shapes:
                # Text frames
                try:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame:
                        # shape.text could be used, but iterating paragraphs/runs gives finer control if needed later
                        for para in shape.text_frame.paragraphs:
                            text = "".join(run.text or "" for run in para.runs) if para.runs else (para.text or "")
                            if _append_and_maybe_stop(chunks, _safe_str(text) + "\n", max_chars):
                                return normalize_text("".join(chunks))
                except Exception:
                    # Ignore malformed shapes and continue
                    pass

                # Tables
                try:
                    if getattr(shape, "has_table", False):
                        tbl = shape.table
                        for r in tbl.rows:
                            for c in r.cells:
                                if _append_and_maybe_stop(chunks, _safe_str(c.text) + "\n", max_chars):
                                    return normalize_text("".join(chunks))
                except Exception:
                    pass

    except Exception:
        # Propagate for processor to handle
        raise

    return normalize_text("".join(chunks))


def extract_xlsx_text(
    path: PathLike,
    max_chars: int,
    include_hidden: bool = True,
    skip_formulas: bool = True,
) -> str:
    """
    Extract plain text from a .xlsx file using openpyxl in read-only mode:
    - Include text from hidden/veryHidden sheets if include_hidden=True.
    - Skip formula cells entirely if skip_formulas=True.
    - Only collect string-like cell values; skip numbers/dates/errors.
    - Stop early once 'max_chars' is reached.
    """
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError("openpyxl is required to extract XLSX text") from exc

    p = Path(path)
    chunks: List[str] = []
    try:
        # read_only=True is memory-friendly; data_only=True returns evaluated values for formulas,
        # but since we skip formulas by data_type check, this is just a safety.
        wb = load_workbook(filename=str(p), read_only=True, data_only=True)

        for ws in wb.worksheets:
            # Sheet visibility handling
            state = getattr(ws, "sheet_state", "visible")
            if not include_hidden and state != "visible":
                continue

            # Iterate cells in streaming mode; need values_only=False so we can check data_type
            for row in ws.iter_rows(values_only=False):
                for cell in row:
                    try:
                        # Skip formulas
                        if skip_formulas and getattr(cell, "data_type", None) == "f":
                            continue

                        val = cell.value
                        # Skip obvious formulas (fallback; shouldn't be needed if data_type == 'f')
                        if skip_formulas and isinstance(val, str) and val.startswith("="):
                            continue

                        # Only collect string-like content
                        if isinstance(val, str) and val.strip():
                            if _append_and_maybe_stop(chunks, val + "\n", max_chars):
                                return normalize_text("".join(chunks))
                    except Exception:
                        # Ignore bad cells and keep going
                        continue

    except Exception:
        # Propagate for processor to handle
        raise

    return normalize_text("".join(chunks))


def extract_pdf_text(path: PathLike, max_chars: int) -> str:
    """
    Extract plain text from a .pdf file using pypdf:
    - Skip encrypted PDFs by raising a PdfReadError (processor will set read_error).
    - For each page, use page.extract_text() (may return None).
    - Stop early once 'max_chars' is reached.

    Note: Text extraction quality depends on how the PDF was created.
    Scanned/image-based PDFs may yield little or no text.
    """
    try:
        from pypdf import PdfReader
        from pypdf.errors import PdfReadError
    except ImportError as exc:
        raise RuntimeError("pypdf is required to extract PDF text") from exc

    p = Path(path)
    chunks: List[str] = []
    try:
        reader = PdfReader(str(p))
        if getattr(reader, "is_encrypted", False):
            # Signal to caller; they will set metadata["read_error"]
            raise PdfReadError("Encrypted PDF: cannot extract text")

        for page in reader.pages:
            try:
                txt = page.extract_text() or ""
                if txt:
                    if _append_and_maybe_stop(chunks, txt + "\n", max_chars):
                        return normalize_text("".join(chunks))
            except Exception:
                # Skip problematic pages
                continue

    except Exception:
        # Propagate for processor to handle
        raise

    return normalize_text("".join(chunks))
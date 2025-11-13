# processors/pptx_processor.py
"""
PPTX technician (FileProcessor): builds a read-only FileArtifact for .pptx files.

Existing metadata produced (consumed by current checks):
- slide_count: int
- total_shapes: int
- has_any_notes: bool
- core_author/core_created/core_modified: core properties
- comments_count: int (counts <p:cm> across /ppt/comments/comment*.xml)
- read_error / read_error_detail on core parse failures

New in Phase 3 (for spelling now, grammar later):
- text_sample: normalized plain text sample gathered from slide text frames and table cells
- text_length: length of text_sample
- token_count: basic English word token count (lowercased)
- text_extraction_error / text_extraction_error_detail: flags if text extraction fails

Design (SOLID):
- SRP: Processor gathers metadata only (including a read-only text sample); checks make policy decisions.
- OCP: Adding new checks reads metadata; orchestrator unchanged.
- LSP: Still honors FileProcessor; orchestrator treats it the same.
- ISP: Checks depend only on FileArtifact (no parsing libraries).
- DIP: Checks depend on abstract text (string in metadata), not python-pptx.

Implementation notes:
- Use python-pptx for high-level properties (slides, shapes, notes, core props).
- Count comments by streaming comment XML parts.
- Use utils.text_extract.extract_pptx_text to build a single text sample
  (early-stopped at max_text_chars from config).
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Optional
from zipfile import ZipFile
from xml.etree.ElementTree import iterparse

from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError

from core.interfaces import FileProcessor
from core.models import FileArtifact
from core.registry import register_processor

# NEW: text extraction helpers and config
from utils.text_extract import extract_pptx_text, tokenize_words
from infra.config_loader import load_config


class PptxProcessor(FileProcessor):
    def supports(self):
        return [".pptx"]

    def build_artifact(self, path: Path) -> FileArtifact:
        size = path.stat().st_size
        metadata: Dict[str, Any] = {
            "kind": "pptx",
            "slide_count": None,
            "total_shapes": None,
            "has_any_notes": None,
            "core_author": None,
            "core_created": None,
            "core_modified": None,
            "comments_count": 0,
            "read_error": False,
            # NEW: text extraction fields (for spelling/grammar checks)
            "text_sample": "",
            "text_length": 0,
            "token_count": 0,
            "text_extraction_error": False,
            # "text_extraction_error_detail": "...",  # set only on error
        }

        # 1) High-level facts via python-pptx (slides, shapes, core props)
        try:
            prs = Presentation(str(path))
            slides = list(prs.slides)
            metadata["slide_count"] = len(slides)
            metadata["total_shapes"] = sum(len(s.shapes) for s in slides)
            try:
                # has_notes_slide is not always reliable across versions, so we guard
                metadata["has_any_notes"] = any(
                    getattr(s, "has_notes_slide", False) for s in slides
                )
            except Exception:
                # Fallback: try creating notes_slide accessor
                def _has_notes(slide):
                    try:
                        return slide.notes_slide is not None
                    except Exception:
                        return False

                metadata["has_any_notes"] = any(_has_notes(s) for s in slides)

            core = prs.core_properties
            metadata["core_author"] = getattr(core, "author", None)
            metadata["core_created"] = _safe_iso(getattr(core, "created", None))
            metadata["core_modified"] = _safe_iso(getattr(core, "modified", None))
        except Exception:
            # Keep going; we can still count comments via ZIP scan and extract text separately.
            pass

        # 2) Count comments via ZIP scan of /ppt/comments/comment*.xml (<p:cm>)
        try:
            with ZipFile(str(path)) as zf:
                total = 0
                for name in zf.namelist():
                    if name.startswith("ppt/comments/comment") and name.endswith(".xml"):
                        total += _count_tag_local(zf, name, "cm")
                metadata["comments_count"] = total
        except (PptxPackageNotFoundError, OSError, ValueError, KeyError, Exception) as exc:
            metadata["read_error"] = True
            metadata["read_error_detail"] = str(exc) or exc.__class__.__name__

        # 3) NEW: Plain-text sample extraction (for spelling/grammar checks)
        # We attempt extraction even if earlier steps encountered issues, so text-only checks can still run.
        try:
            cfg = load_config()
            if cfg.get("enable_spelling", True) or cfg.get("enable_grammar", False):
                max_chars = int(cfg.get("max_text_chars", 5_000_000))
                sample = extract_pptx_text(path, max_chars)
                metadata["text_sample"] = sample
                metadata["text_length"] = len(sample)
                metadata["token_count"] = len(tokenize_words(sample))
        except Exception as exc:
            metadata["text_sample"] = ""
            metadata["text_length"] = 0
            metadata["token_count"] = 0
            metadata["text_extraction_error"] = True
            metadata["text_extraction_error_detail"] = str(exc) or exc.__class__.__name__

        return FileArtifact(
            path=path,
            extension=".pptx",
            size_bytes=size,
            metadata=metadata,
        )


# --- helpers ---

def _safe_iso(dt) -> Optional[str]:
    try:
        return dt.isoformat() if dt is not None else None
    except Exception:
        return None


def _count_tag_local(zf: ZipFile, member: str, local_name: str) -> int:
    """
    Count how many elements with the given local-name appear in an XML part.
    For PPTX comments we look for local-name 'cm' inside comment parts.
    """
    count = 0
    with zf.open(member) as fp:
        for _event, elem in iterparse(fp, events=("end",)):
            tag = elem.tag
            if "}" in tag:
                tag = tag.split("}", 1)[1]
            if tag == local_name:
                count += 1
            elem.clear()
    return count


# Register on import
register_processor(PptxProcessor())